#!/usr/bin/env python3
"""
LoBRA Data Pipeline - File Preprocessor
Converts various file formats into markdown for ingestion
"""

import os
import sys
import yaml
import shutil
from pathlib import Path
from datetime import datetime
from tqdm import tqdm
import mimetypes

# Initialize mimetypes
mimetypes.init()

def log_info(msg):
    print(f"[INFO] {msg}")

def log_success(msg):
    print(f"[SUCCESS] {msg}")

def log_warning(msg):
    print(f"[WARNING] {msg}")

def log_error(msg):
    print(f"[ERROR] {msg}", file=sys.stderr)


class FileConverter:
    """Base converter class"""
    
    def __init__(self, config):
        self.config = config
        self.supported_formats = []
        
    def can_convert(self, file_path):
        """Check if this converter can handle the file"""
        return file_path.suffix.lower() in self.supported_formats
    
    def convert(self, input_path, output_path):
        """Convert file to markdown"""
        raise NotImplementedError


class WordConverter(FileConverter):
    """Convert Word documents to markdown"""
    
    def __init__(self, config):
        super().__init__(config)
        self.supported_formats = ['.docx', '.doc']
        
    def convert(self, input_path, output_path):
        try:
            import docx2txt
            
            # Extract text from docx
            text = docx2txt.process(str(input_path))
            
            # Create markdown with metadata
            metadata = self._create_metadata(input_path)
            markdown = self._format_markdown(metadata, text)
            
            output_path.write_text(markdown, encoding='utf-8')
            return True
            
        except ImportError:
            log_warning("docx2txt not installed. Install with: pip install docx2txt")
            return False
        except Exception as e:
            log_error(f"Failed to convert {input_path.name}: {e}")
            return False
    
    def _create_metadata(self, input_path):
        stat = input_path.stat()
        return {
            'title': input_path.stem,
            'date': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d'),
            'source': str(input_path),
            'format': 'Word Document',
            'converted': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
    
    def _format_markdown(self, metadata, text):
        md = "---\n"
        for key, value in metadata.items():
            md += f"{key}: \"{value}\"\n"
        md += "---\n\n"
        md += text
        return md


class PowerPointConverter(FileConverter):
    """Convert PowerPoint presentations to markdown"""
    
    def __init__(self, config):
        super().__init__(config)
        self.supported_formats = ['.pptx', '.ppt']
        
    def convert(self, input_path, output_path):
        try:
            from pptx import Presentation
            
            prs = Presentation(str(input_path))
            text_runs = []
            
            for i, slide in enumerate(prs.slides, 1):
                text_runs.append(f"\n## Slide {i}\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_runs.append(shape.text)
                        text_runs.append("\n")
            
            text = "\n".join(text_runs)
            
            # Create markdown with metadata
            metadata = self._create_metadata(input_path, len(prs.slides))
            markdown = self._format_markdown(metadata, text)
            
            output_path.write_text(markdown, encoding='utf-8')
            return True
            
        except ImportError:
            log_warning("python-pptx not installed. Install with: pip install python-pptx")
            return False
        except Exception as e:
            log_error(f"Failed to convert {input_path.name}: {e}")
            return False
    
    def _create_metadata(self, input_path, slide_count):
        stat = input_path.stat()
        return {
            'title': input_path.stem,
            'date': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d'),
            'source': str(input_path),
            'format': 'PowerPoint Presentation',
            'slides': str(slide_count),
            'converted': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
    
    def _format_markdown(self, metadata, text):
        md = "---\n"
        for key, value in metadata.items():
            md += f"{key}: \"{value}\"\n"
        md += "---\n\n"
        md += f"# {metadata['title']}\n\n"
        md += text
        return md


class ExcelConverter(FileConverter):
    """Convert Excel spreadsheets to markdown tables"""
    
    def __init__(self, config):
        super().__init__(config)
        self.supported_formats = ['.xlsx', '.xls', '.csv']
        
    def convert(self, input_path, output_path):
        try:
            import pandas as pd
            
            # Read based on file type
            if input_path.suffix.lower() == '.csv':
                df = pd.read_csv(input_path)
            else:
                df = pd.read_excel(input_path, sheet_name=None)  # Read all sheets
            
            # Convert to markdown
            if isinstance(df, dict):  # Multiple sheets
                text = self._convert_multiple_sheets(df)
            else:  # Single sheet or CSV
                text = f"## Data\n\n{df.to_markdown(index=False)}"
            
            # Create markdown with metadata
            metadata = self._create_metadata(input_path)
            markdown = self._format_markdown(metadata, text)
            
            output_path.write_text(markdown, encoding='utf-8')
            return True
            
        except ImportError:
            log_warning("pandas not installed. Install with: pip install pandas openpyxl tabulate")
            return False
        except Exception as e:
            log_error(f"Failed to convert {input_path.name}: {e}")
            return False
    
    def _convert_multiple_sheets(self, sheets_dict):
        text = ""
        for sheet_name, df in sheets_dict.items():
            text += f"\n## Sheet: {sheet_name}\n\n"
            text += df.to_markdown(index=False)
            text += "\n"
        return text
    
    def _create_metadata(self, input_path):
        stat = input_path.stat()
        return {
            'title': input_path.stem,
            'date': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d'),
            'source': str(input_path),
            'format': 'Spreadsheet',
            'converted': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
    
    def _format_markdown(self, metadata, text):
        md = "---\n"
        for key, value in metadata.items():
            md += f"{key}: \"{value}\"\n"
        md += "---\n\n"
        md += f"# {metadata['title']}\n\n"
        md += text
        return md


class HTMLConverter(FileConverter):
    """Convert HTML files to markdown"""
    
    def __init__(self, config):
        super().__init__(config)
        self.supported_formats = ['.html', '.htm']
        
    def convert(self, input_path, output_path):
        try:
            from bs4 import BeautifulSoup
            import html2text
            
            # Read HTML
            html_content = input_path.read_text(encoding='utf-8')
            
            # Convert to markdown
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = False
            text = h.handle(html_content)
            
            # Create markdown with metadata
            metadata = self._create_metadata(input_path)
            markdown = self._format_markdown(metadata, text)
            
            output_path.write_text(markdown, encoding='utf-8')
            return True
            
        except ImportError:
            log_warning("html2text not installed. Install with: pip install html2text beautifulsoup4")
            return False
        except Exception as e:
            log_error(f"Failed to convert {input_path.name}: {e}")
            return False
    
    def _create_metadata(self, input_path):
        stat = input_path.stat()
        return {
            'title': input_path.stem,
            'date': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d'),
            'source': str(input_path),
            'format': 'HTML',
            'converted': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
    
    def _format_markdown(self, metadata, text):
        md = "---\n"
        for key, value in metadata.items():
            md += f"{key}: \"{value}\"\n"
        md += "---\n\n"
        md += text
        return md


class ImageConverter(FileConverter):
    """Extract text from images using OCR"""
    
    def __init__(self, config):
        super().__init__(config)
        self.supported_formats = ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']
        
    def convert(self, input_path, output_path):
        try:
            from PIL import Image
            import pytesseract
            
            # Open image and extract text
            img = Image.open(input_path)
            text = pytesseract.image_to_string(img)
            
            if not text.strip():
                log_warning(f"No text found in {input_path.name}")
                return False
            
            # Create markdown with metadata
            metadata = self._create_metadata(input_path)
            markdown = self._format_markdown(metadata, text)
            
            output_path.write_text(markdown, encoding='utf-8')
            return True
            
        except ImportError:
            log_warning("pytesseract or PIL not installed. Install with: pip install pytesseract Pillow")
            log_warning("Also install tesseract: brew install tesseract (macOS) or apt install tesseract-ocr (Linux)")
            return False
        except Exception as e:
            log_error(f"Failed to convert {input_path.name}: {e}")
            return False
    
    def _create_metadata(self, input_path):
        stat = input_path.stat()
        return {
            'title': input_path.stem,
            'date': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d'),
            'source': str(input_path),
            'format': 'Image (OCR)',
            'converted': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
    
    def _format_markdown(self, metadata, text):
        md = "---\n"
        for key, value in metadata.items():
            md += f"{key}: \"{value}\"\n"
        md += "---\n\n"
        md += f"# {metadata['title']}\n\n"
        md += "**Note:** Text extracted from image using OCR\n\n"
        md += text
        return md


class PDFConverter(FileConverter):
    """Convert PDF files to markdown using PyPDF"""
    
    def __init__(self, config):
        super().__init__(config)
        self.supported_formats = ['.pdf']
        
    def convert(self, input_path, output_path):
        try:
            from pypdf import PdfReader
            
            reader = PdfReader(str(input_path))
            text_parts = []
            
            for i, page in enumerate(reader.pages, 1):
                page_text = page.extract_text()
                if page_text.strip():
                    text_parts.append(f"\n## Page {i}\n\n{page_text}")
            
            text = "\n".join(text_parts)
            
            if not text.strip():
                log_warning(f"No text extracted from {input_path.name}")
                return False
            
            # Create markdown with metadata
            metadata = self._create_metadata(input_path, len(reader.pages))
            markdown = self._format_markdown(metadata, text)
            
            output_path.write_text(markdown, encoding='utf-8')
            return True
            
        except ImportError:
            log_warning("pypdf not installed. Install with: pip install pypdf")
            return False
        except Exception as e:
            log_error(f"Failed to convert {input_path.name}: {e}")
            return False
    
    def _create_metadata(self, input_path, page_count):
        stat = input_path.stat()
        return {
            'title': input_path.stem,
            'date': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d'),
            'source': str(input_path),
            'format': 'PDF Document',
            'pages': str(page_count),
            'converted': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
    
    def _format_markdown(self, metadata, text):
        md = "---\n"
        for key, value in metadata.items():
            md += f"{key}: \"{value}\"\n"
        md += "---\n\n"
        md += f"# {metadata['title']}\n\n"
        if metadata.get('pages'):
            md += f"**Pages:** {metadata['pages']}\n\n"
        md += text
        return md


class MarkdownConverter(FileConverter):
    """Copy Markdown files with metadata front-matter"""
    
    def __init__(self, config):
        super().__init__(config)
        self.supported_formats = ['.md', '.markdown']
        
    def convert(self, input_path, output_path):
        try:
            import frontmatter
            
            # Try to load existing front-matter
            try:
                post = frontmatter.load(input_path)
                content = post.content
                existing_meta = dict(post.metadata)
            except:
                # No front-matter, use entire file as content
                content = input_path.read_text(encoding='utf-8')
                existing_meta = {}
            
            # Create/update metadata
            stat = input_path.stat()
            metadata = {
                'title': existing_meta.get('title', input_path.stem),
                'date': existing_meta.get('date', datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d')),
                'source': existing_meta.get('source', str(input_path)),
                'format': 'Markdown',
                'processed': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            }
            
            # Merge existing metadata (preserve custom fields)
            for key, value in existing_meta.items():
                if key not in ['title', 'date', 'source', 'format']:  # Don't override core fields
                    metadata[key] = value
            
            # Format with front-matter
            md = "---\n"
            for key, value in metadata.items():
                if isinstance(value, list):
                    md += f"{key}: {value}\n"
                elif isinstance(value, (int, float, bool)):
                    md += f"{key}: {value}\n"
                else:
                    md += f"{key}: \"{value}\"\n"
            md += "---\n\n"
            md += content
            
            output_path.write_text(md, encoding='utf-8')
            return True
            
        except ImportError:
            log_warning("python-frontmatter not installed. Install with: pip install python-frontmatter")
            # Fallback: simple copy with basic metadata
            return self._simple_copy(input_path, output_path)
        except Exception as e:
            log_error(f"Failed to convert {input_path.name}: {e}")
            # Fallback: simple copy
            return self._simple_copy(input_path, output_path)
    
    def _simple_copy(self, input_path, output_path):
        """Fallback: simple copy with basic metadata"""
        try:
            content = input_path.read_text(encoding='utf-8')
            stat = input_path.stat()
            
            metadata = {
                'title': input_path.stem,
                'date': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d'),
                'source': str(input_path),
                'format': 'Markdown',
                'processed': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            }
            
            md = "---\n"
            for key, value in metadata.items():
                md += f"{key}: \"{value}\"\n"
            md += "---\n\n"
            md += content
            
            output_path.write_text(md, encoding='utf-8')
            return True
        except Exception as e:
            log_error(f"Failed to copy {input_path.name}: {e}")
            return False


class EPUBConverter(FileConverter):
    """Convert EPUB books to markdown"""
    
    def __init__(self, config):
        super().__init__(config)
        self.supported_formats = ['.epub']
        
    def convert(self, input_path, output_path):
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            
            book = epub.read_epub(str(input_path))
            chapters = []
            
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text = soup.get_text()
                    if text.strip():
                        chapters.append(text)
            
            text = "\n\n---\n\n".join(chapters)
            
            # Create markdown with metadata
            metadata = self._create_metadata(input_path, book)
            markdown = self._format_markdown(metadata, text)
            
            output_path.write_text(markdown, encoding='utf-8')
            return True
            
        except ImportError:
            log_warning("ebooklib not installed. Install with: pip install ebooklib beautifulsoup4")
            return False
        except Exception as e:
            log_error(f"Failed to convert {input_path.name}: {e}")
            return False
    
    def _create_metadata(self, input_path, book):
        stat = input_path.stat()
        
        # Try to extract book metadata
        title = input_path.stem
        author = "Unknown"
        
        try:
            title = book.get_metadata('DC', 'title')[0][0] if book.get_metadata('DC', 'title') else title
            author = book.get_metadata('DC', 'creator')[0][0] if book.get_metadata('DC', 'creator') else author
        except:
            pass
        
        return {
            'title': title,
            'author': author,
            'date': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d'),
            'source': str(input_path),
            'format': 'EPUB Book',
            'converted': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
    
    def _format_markdown(self, metadata, text):
        md = "---\n"
        for key, value in metadata.items():
            md += f"{key}: \"{value}\"\n"
        md += "---\n\n"
        md += f"# {metadata['title']}\n\n"
        if metadata.get('author') != "Unknown":
            md += f"**Author:** {metadata['author']}\n\n"
        md += text
        return md


class DataPipeline:
    """Main data pipeline orchestrator"""
    
    def __init__(self, config_path="config.yaml"):
        self.config = yaml.safe_load(open(config_path))
        self.input_dir = Path("./inbox")  # New files go here
        self.output_dir = Path(self.config['vault_path'])
        self.processed_dir = Path("./processed")  # Originals moved here
        
        # Initialize converters
        # Order matters: Markdown first (fastest check), then others
        self.converters = [
            MarkdownConverter(self.config),  # Fast copy with metadata
            PDFConverter(self.config),        # PDF text extraction
            WordConverter(self.config),
            PowerPointConverter(self.config),
            ExcelConverter(self.config),
            HTMLConverter(self.config),
            ImageConverter(self.config),
            EPUBConverter(self.config),
        ]
        
        # Create directories
        self.input_dir.mkdir(exist_ok=True)
        self.processed_dir.mkdir(exist_ok=True)
        self.output_dir.mkdir(exist_ok=True)
    
    def process_file(self, file_path):
        """Process a single file"""
        log_info(f"Processing: {file_path.name}")
        
        # Find appropriate converter
        converter = None
        for conv in self.converters:
            if conv.can_convert(file_path):
                converter = conv
                break
        
        if not converter:
            log_warning(f"No converter found for {file_path.name}")
            return False
        
        # Generate output path
        output_path = self.output_dir / f"{file_path.stem}.md"
        
        # Handle filename conflicts
        counter = 1
        while output_path.exists():
            output_path = self.output_dir / f"{file_path.stem}_{counter}.md"
            counter += 1
        
        # Convert
        if converter.convert(file_path, output_path):
            log_success(f"Converted to: {output_path.name}")
            
            # Move original to processed directory
            processed_path = self.processed_dir / file_path.name
            counter = 1
            while processed_path.exists():
                processed_path = self.processed_dir / f"{file_path.stem}_{counter}{file_path.suffix}"
                counter += 1
            
            shutil.move(str(file_path), str(processed_path))
            log_info(f"Original moved to: {processed_path}")
            return True
        
        return False
    
    def process_all(self):
        """Process all files in inbox"""
        files = list(self.input_dir.rglob("*"))  # rglob to handle subdirectories
        files = [f for f in files if f.is_file() and not f.name.startswith('.')]
        
        if not files:
            log_info("No files found in inbox/")
            print("\nTo add files for processing:")
            print(f"  cp your-file.docx {self.input_dir}/")
            print(f"  python scripts/preprocess.py")
            return
        
        log_info(f"Found {len(files)} file(s) to process")
        
        success_count = 0
        failed_count = 0
        
        for file_path in tqdm(files, desc="Processing files"):
            if self.process_file(file_path):
                success_count += 1
            else:
                failed_count += 1
        
        print("\n" + "="*50)
        log_success(f"Processed: {success_count} files")
        if failed_count > 0:
            log_warning(f"Failed: {failed_count} files")
        print("="*50)
        
        if success_count > 0:
            print("\nNext steps:")
            print("  1. Review converted files in vault/")
            print("  2. Run: make ingest")
            print("  3. Query your knowledge: make ask q=\"your question\"")


def main():
    print("╔═══════════════════════════════════════════════════════╗")
    print("║         LoBRA Data Pipeline - Preprocessor           ║")
    print("╚═══════════════════════════════════════════════════════╝")
    print()
    
    try:
        pipeline = DataPipeline()
        pipeline.process_all()
    except Exception as e:
        log_error(f"Pipeline failed: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()

